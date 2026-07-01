"""scripts/modify_rag_ppt.py
기존 "RAG 부분.pptx" → "RAG_부분_v2.pptx" 생성
ppt_plan_v2.md 기반, 5슬라이드 자유 양식:
  Slide 14: Vector DB 구성
  Slide 15: DB 구축 현황
  Slide 16: 질문 분석 — 의도 분류 & 엔티티 추출
  Slide 17: 검색 과정 — Hybrid Search & Re-ranking
  Slide 18: 검색 예시 — 실제 동작
"""

import copy
import os
import sys

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

INPUT   = "RAG 부분.pptx"
OUTPUT  = "RAG_부분_v2.pptx"
IMG_DIR = "image"

# ── 색상 ──────────────────────────────────────────────────────────────────
NAVY   = "082F76"
TEAL   = "25577C"
BLUE   = "0070C0"
LBLUE  = "D6E4F0"
DARK   = "2D2D2D"
LGRAY  = "E0E0E0"
ORANGE = "E86A2C"
PGRAY  = "F2F2F2"
PGRAY2 = "A0A0A0"
WHITE  = "FFFFFF"
GRAY20 = "333333"
GRAY50 = "808080"
GREEN  = "1F6B2E"
LGREEN = "E2F0E6"

# 5개 테이블 색상
TABLE_COLORS = {
    "accidents": ("C0392B", "FADBD8"),  # 빨강
    "chemicals": ("8E44AD", "E8DAEF"),  # 보라
    "laws":      ("1A5276", "D6EAF8"),  # 파랑
    "designs":   ("1E8449", "D5F5E3"),  # 초록
    "basics":    ("D35400", "FAE5D3"),  # 주황
}

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

def aq(tag): return f"{{{NS_A}}}{tag}"
def pq(tag): return f"{{{NS_P}}}{tag}"

# ── 슬라이드 복제 / 재정렬 ────────────────────────────────────────────────

def duplicate_slide(prs, src_idx):
    src = prs.slides[src_idx]
    new_slide = prs.slides.add_slide(prs.slide_layouts[-1])
    spTree = new_slide.shapes._spTree
    for sp in list(spTree):
        spTree.remove(sp)
    for sp in src.shapes._spTree:
        spTree.append(copy.deepcopy(sp))
    return new_slide


def reorder_slides(prs, order):
    sldIdLst = prs.slides._sldIdLst
    elems = list(sldIdLst)
    for child in list(sldIdLst):
        sldIdLst.remove(child)
    for i in order:
        sldIdLst.append(elems[i])


# ── 배경 초기화: 제목/타이틀바만 남기고 나머지 제거 ────────────────────────

KEEP_BG = {"TextBox 1", "TextBox 2", "직사각형 10"}

def clear_slide(slide):
    """배경 장식(제목·타이틀바) 제외 모두 제거 — 자유 양식을 위해 패널 헤더도 제거"""
    for shape in list(slide.shapes):
        if shape.name not in KEEP_BG:
            shape._element.getparent().remove(shape._element)


# ── 배경 텍스트 업데이트 ───────────────────────────────────────────────────

def _set_first_run(para_el, text):
    runs = para_el.findall(aq("r"))
    if runs:
        t = runs[0].find(aq("t"))
        if t is not None:
            t.text = text
        for r in runs[1:]:
            para_el.remove(r)


def update_title(slide, text):
    for shape in slide.shapes:
        if shape.name == "TextBox 1":
            paras = shape._element.findall(f".//{aq('p')}")
            if paras:
                _set_first_run(paras[0], text)
            break


def update_titlebar(slide, line1, line2=""):
    for shape in slide.shapes:
        if shape.name == "직사각형 10":
            paras = shape._element.findall(f".//{aq('p')}")
            if len(paras) >= 1:
                _set_first_run(paras[0], line1)
            if len(paras) >= 2:
                _set_first_run(paras[1], line2)
            break


# ── 도형 추가 헬퍼 ────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_hex=None, line_hex=None, line_pt=0):
    shape = slide.shapes.add_shape(1, Cm(l), Cm(t), Cm(w), Cm(h))
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    else:
        shape.fill.background()
    if line_hex and line_pt > 0:
        shape.line.color.rgb = RGBColor.from_string(line_hex)
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, l, t, w, h, lines, word_wrap=True):
    box = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf  = box.text_frame
    tf.word_wrap = word_wrap
    for i, line in enumerate(lines):
        if isinstance(line, str):
            d = {"text": line}
        else:
            d = line
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = d.get("align", PP_ALIGN.LEFT)
        run = para.add_run()
        run.text = d["text"]
        run.font.size = Pt(d.get("sz", 9.5))
        run.font.bold = d.get("bold", False)
        if d.get("color_hex"):
            run.font.color.rgb = RGBColor.from_string(d["color_hex"])
        if d.get("font"):
            try:
                run.font._element.get_or_add_latin().set("typeface", d["font"])
                run.font._element.get_or_add_ea().set("typeface", d["font"])
            except Exception:
                pass
    return box


def add_code_box(slide, l, t, w, h, lines, title=None, font_sz=8):
    bg = add_rect(slide, l, t, w, h, fill_hex=DARK)
    tf = bg.text_frame
    tf.word_wrap = False
    first = True
    if title:
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size = Pt(8.5)
        r.font.bold = True
        r.font.color.rgb = RGBColor.from_string(ORANGE)
        try:
            r.font._element.get_or_add_latin().set("typeface", "Consolas")
        except Exception:
            pass
        first = False
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = line
        r.font.size  = Pt(font_sz)
        r.font.color.rgb = RGBColor.from_string(LGRAY)
        try:
            r.font._element.get_or_add_latin().set("typeface", "Consolas")
        except Exception:
            pass
    return bg


def add_placeholder(slide, l, t, w, h, caption):
    box = add_rect(slide, l, t, w, h, fill_hex=PGRAY, line_hex=PGRAY2, line_pt=1.2)
    tf  = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = "[ 📷  스크린샷 ]"
    r1.font.size  = Pt(10)
    r1.font.bold  = True
    r1.font.color.rgb = RGBColor.from_string(GRAY50)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = caption
    r2.font.size  = Pt(8.5)
    r2.font.color.rgb = RGBColor.from_string(GRAY50)
    return box


def add_image_or_placeholder(slide, img_name, l, t, w, h, caption):
    img_path = os.path.join(IMG_DIR, img_name)
    if os.path.isfile(img_path):
        slide.shapes.add_picture(img_path, Cm(l), Cm(t), Cm(w), Cm(h))
        print(f"  이미지 삽입: {img_name}")
    else:
        add_placeholder(slide, l, t, w, h, caption)
        print(f"  플레이스홀더: {img_name}")


def section_header(slide, l, t, w, text, color_hex=BLUE):
    """섹션 소제목 — 왼쪽 색상 바 + 텍스트"""
    add_rect(slide, l, t, 0.25, 0.55, fill_hex=color_hex)
    add_textbox(slide, l + 0.35, t, w - 0.35, 0.55, [
        {"text": text, "sz": 10.5, "bold": True, "color_hex": color_hex}
    ])


# ── 표 헬퍼 ──────────────────────────────────────────────────────────────

def _set_cell_fill(cell, hex_color):
    tc = cell._tc
    tcPr = tc.find(aq("tcPr"))
    if tcPr is None:
        tcPr = etree.SubElement(tc, aq("tcPr"))
    for child in list(tcPr):
        if child.tag in (aq("solidFill"), aq("gradFill"), aq("noFill")):
            tcPr.remove(child)
    sf = etree.SubElement(tcPr, aq("solidFill"))
    etree.SubElement(sf, aq("srgbClr")).set("val", hex_color)


def add_table(slide, data, l, t, w, h,
              col_widths=None,
              header_fill=NAVY, header_text=WHITE,
              first_col_fill=LBLUE, first_col_text=NAVY,
              even_fill="F5F9FC", odd_fill=WHITE,
              font_sz=8.5, header_sz=9.5,
              h_header=PP_ALIGN.CENTER,
              h_data=PP_ALIGN.LEFT,
              h_first=PP_ALIGN.LEFT):
    nrows = len(data)
    ncols = max(len(row) for row in data)
    gf  = slide.shapes.add_table(nrows, ncols, Cm(l), Cm(t), Cm(w), Cm(h))
    tbl = gf.table
    if col_widths:
        total = sum(col_widths)
        total_emu = Cm(w)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(total_emu * cw / total)
    for r in range(nrows):
        for c in range(ncols):
            cell = tbl.cell(r, c)
            text = data[r][c] if c < len(data[r]) else ""
            tf   = cell.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            if r == 0:
                para.alignment = h_header
                _set_cell_fill(cell, header_fill)
                run = para.add_run()
                run.text = text
                run.font.size  = Pt(header_sz)
                run.font.bold  = True
                run.font.color.rgb = RGBColor.from_string(header_text)
            elif c == 0:
                para.alignment = h_first
                _set_cell_fill(cell, first_col_fill)
                run = para.add_run()
                run.text = text
                run.font.size  = Pt(font_sz)
                run.font.bold  = True
                run.font.color.rgb = RGBColor.from_string(first_col_text)
            else:
                para.alignment = h_data
                _set_cell_fill(cell, even_fill if r % 2 == 0 else odd_fill)
                run = para.add_run()
                run.text = text
                run.font.size  = Pt(font_sz)
                run.font.color.rgb = RGBColor.from_string(GRAY20)
    return gf


# ════════════════════════════════════════════════════════════════════════════
# Slide 1: DB 구축 — 전체 흐름
# ════════════════════════════════════════════════════════════════════════════

def build_slide14(slide):
    clear_slide(slide)
    update_title(slide, "RAG 구축 - DB 구축 흐름")
    update_titlebar(slide,
        "DB 구축 흐름",
        "5종 전문 문서를 파싱·청킹·임베딩하여 LanceDB에 저장")

    Y = 6.3

    # === 파이프라인 도식 ===
    section_header(slide, 0.8, Y, 32, "데이터 → 벡터DB 구축 파이프라인")
    Y += 0.65

    BW = 5.9
    AW = 0.6
    BH = 2.1
    X  = 0.8

    steps = [
        ("원본 파일\n(PDF / Excel)", PGRAY, PGRAY2, GRAY20),
        ("파싱\n(pdfplumber /\nopenpyxl)", PGRAY, PGRAY2, GRAY20),
        ("청킹\n(테이블별\n방식 상이)", LBLUE, BLUE, NAVY),
        ("BGE-M3\n임베딩\n(→ 1,024차원\n벡터)", LBLUE, BLUE, NAVY),
        ("LanceDB\n저장\n(로컬 파일 기반\n서버 불필요)", LGREEN, GREEN, GREEN),
    ]

    for i, (label, fill, border, text_c) in enumerate(steps):
        add_rect(slide, X, Y, BW, BH, fill_hex=fill, line_hex=border, line_pt=1.0)
        add_textbox(slide, X + 0.1, Y + 0.2, BW - 0.2, BH - 0.2, [
            {"text": label, "sz": 9, "bold": True, "color_hex": text_c,
             "align": PP_ALIGN.CENTER}
        ])
        X += BW
        if i < 4:
            add_textbox(slide, X, Y + BH/2 - 0.22, AW, 0.44, [
                {"text": "→", "sz": 13, "bold": True, "color_hex": GRAY20,
                 "align": PP_ALIGN.CENTER}
            ])
            X += AW

    Y += BH + 0.3
    add_textbox(slide, 0.8, Y, 32.0, 0.7, [
        {"text": "▸  청킹 방식은 테이블마다 다름: 행 단위(accidents) / 페이지 단위(chemicals) / 조항 단위(laws) / 800자 슬라이딩(designs, basics)",
         "sz": 8.5, "color_hex": GRAY50},
        {"text": "▸  임베딩 모델: BAAI/bge-m3  (한국어·영어 다국어, 1,024차원)",
         "sz": 8.5, "color_hex": GRAY50},
    ])
    Y += 0.9

    # === 구축 결과 요약 ===
    section_header(slide, 0.8, Y, 32, "구축 결과")
    Y += 0.65

    summary_data = [
        ["테이블",    "원본 파일",  "원본 규모",    "청크 수"],
        ["accidents", "xlsx × 3",  "1,348행",      "2,512건"],
        ["chemicals", "pdf × 1",   "501페이지",    "501건"],
        ["laws",      "pdf × 43",  "2,099페이지",  "3,091건"],
        ["designs",   "pdf × 61",  "1,983페이지",  "3,121건"],
        ["basics",    "pdf × 15",  "8,165페이지",  "29,757건"],
        ["합계",      "",          "",             "38,982건"],
    ]
    add_table(slide, summary_data,
              l=0.8, t=Y, w=32.0, h=5.0,
              col_widths=[4, 4, 6, 4],
              font_sz=9, header_sz=10,
              h_data=PP_ALIGN.CENTER,
              h_first=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# Slide 2: DB 구축 — 상세
# ════════════════════════════════════════════════════════════════════════════

def build_slide14b(slide):
    clear_slide(slide)
    update_title(slide, "RAG 구축 - DB 구축 상세")
    update_titlebar(slide,
        "DB 구축 상세",
        "테이블별 최적화된 청킹 전략 + 핵심 필드별 추가 임베딩")

    Y = 6.3

    # === 청킹 방식 ===
    section_header(slide, 0.8, Y, 32, "테이블별 청킹 방식")
    Y += 0.65

    chunk_data = [
        ["테이블",    "청킹 방식",                        "청킹 이유"],
        ["accidents", "행 단위",                          "사고 1건 = 1행 (물질·설비·원인 통합) — 분리 시 맥락 손실"],
        ["chemicals", "페이지 단위",                      "물질 1종 = 1페이지 단위로 정보 구성됨"],
        ["laws",      "조항 단위",                        "조항 번호+내용을 함께 유지해야 정확한 법적 근거 반환 가능"],
        ["designs",   "800자 슬라이딩 (overlap 100자)",   "지침 본문은 연속 텍스트 — 고정 크기로 균등 분할"],
        ["basics",    "800자 슬라이딩 (overlap 100자)",   "이론 설명은 연속 텍스트 — 고정 크기로 균등 분할"],
    ]
    add_table(slide, chunk_data,
              l=0.8, t=Y, w=32.0, h=3.2,
              col_widths=[3, 5, 12],
              font_sz=8.5, header_sz=9.5,
              h_data=PP_ALIGN.LEFT)
    Y += 3.4

    # === 필드별 추가 임베딩 ===
    section_header(slide, 0.8, Y, 32, "정밀 검색을 위한 필드별 추가 임베딩")
    Y += 0.65

    add_textbox(slide, 0.8, Y, 32.0, 0.9, [
        {"text": "▸  각 청크는 텍스트 전체를 벡터로 변환 (text_vector) — 모든 테이블 공통",
         "sz": 8.5, "color_hex": GRAY20},
        {"text": "▸  추가로, 물질명·설비명·조항번호 같은 핵심 필드를 별도 벡터로 저장",
         "sz": 8.5, "color_hex": GRAY20},
        {"text": '    → 예: "톨루엔 관련 사고만" 검색할 때 material_vector만 따로 검색 → 더 정밀한 결과',
         "sz": 8, "color_hex": GRAY50},
    ])
    Y += 1.0

    vec_data = [
        ["테이블",    "저장되는 벡터",                                       "용도"],
        ["accidents", "text_vector, material_vector, equipment_vector",     "물질·설비 정밀 매칭"],
        ["chemicals", "text_vector, chemical_name_vector",                  "물질명 전용 검색"],
        ["laws",      "text_vector, title_vector, article_vector",          "법령명 + 조항 검색"],
        ["designs",   "text_vector, title_vector, section_vector",          "지침명 + 섹션 검색"],
        ["basics",    "text_vector, title_vector, chapter_vector",          "제목 + 챕터 검색"],
    ]
    add_table(slide, vec_data,
              l=0.8, t=Y, w=32.0, h=3.0,
              col_widths=[2.5, 9.5, 4.5],
              font_sz=8, header_sz=9,
              h_data=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# Slide 15: DB 구축 현황
# ════════════════════════════════════════════════════════════════════════════

def build_slide15(slide):
    clear_slide(slide)
    update_title(slide, "RAG 구축 - DB 구축 현황")
    update_titlebar(slide,
        "DB 구축 현황",
        "로컬 환경에 실제 구축된 Corpus·Vector DB — 폴더 구조 및 실행 로그")

    MID = 16.4   # 좌/우 경계
    Y   = 6.3

    # ── 왼쪽: 폴더 구조 ──
    section_header(slide, 0.8, Y, MID - 1.5, "폴더 구조")
    Y_L = Y + 0.65
    add_code_box(slide, 0.8, Y_L, MID - 1.0, 5.6, [
        "data/",
        "├── corpus/                   ← 원본 PDF / Excel 파일",
        "│   ├── accidents/            xlsx 3개",
        "│   ├── chemicals/            pdf 1개",
        "│   ├── laws/                 pdf 43개",
        "│   ├── designs/              pdf 61개",
        "│   └── basics/               pdf 15개",
        "└── vector_db/                ← LanceDB 임베딩 저장소",
        "    ├── accidents.lance/",
        "    ├── chemicals.lance/",
        "    ├── laws.lance/",
        "    ├── designs.lance/",
        "    └── basics.lance/",
    ], title="▸ 로컬 파일 구조  (LanceDB = 파일 기반 DB, 별도 서버 불필요)")

    add_image_or_placeholder(
        slide, "vector_db_folder.png",
        l=0.8, t=Y_L + 5.75, w=MID - 1.0, h=4.8,
        caption="VS Code 탐색기 — data/vector_db/ 펼치기\n(5개 .lance 폴더 확인)")

    # ── 오른쪽: 실행 화면 ──
    X_R = MID + 0.5
    W_R = 32.8 - X_R
    section_header(slide, X_R, Y, W_R, "실행 화면")

    add_textbox(slide, X_R, Y + 0.65, W_R, 0.45, [
        {"text": "① 벡터 DB 구축 로그  (corpus_loader.py)",
         "sz": 9.5, "bold": True, "color_hex": TEAL}
    ])
    add_image_or_placeholder(
        slide, "corpus_loader_log.png",
        l=X_R, t=Y + 1.2, w=W_R, h=5.1,
        caption="python src/services/corpus_loader.py 실행 마지막 부분\n(capture_guide_slide2.txt 참조)")

    add_textbox(slide, X_R, Y + 6.45, W_R, 0.45, [
        {"text": "② DB 건수 확인  (show_db_status.py)",
         "sz": 9.5, "bold": True, "color_hex": TEAL}
    ])
    add_image_or_placeholder(
        slide, "db_status.png",
        l=X_R, t=Y + 7.0, w=W_R, h=4.0,
        caption="python data/show_db_status.py 전체 출력 캡처")


# ════════════════════════════════════════════════════════════════════════════
# Slide 4: 검색 파이프라인 — 전체 흐름
# ════════════════════════════════════════════════════════════════════════════

def build_slide_pipeline(slide):
    clear_slide(slide)
    update_title(slide, "RAG 구축 - 검색 파이프라인")
    update_titlebar(slide,
        "검색 파이프라인",
        "사용자 질문이 들어오면 어떤 단계를 거쳐 관련 문서가 반환되는가")

    Y = 6.3

    # === 전체 흐름 도식 ===
    section_header(slide, 0.8, Y, 32, "전체 흐름")
    Y += 0.65

    H_ROW  = 0.82
    V_GAP  = 0.1
    TOTAL_H = 4 * H_ROW + 3 * V_GAP

    def row_y(n): return Y + n * (H_ROW + V_GAP)

    # 사용자 질문 박스
    QX, QW = 0.8, 2.2
    add_rect(slide, QX, Y, QW, TOTAL_H, fill_hex="EBF3FB", line_hex=NAVY, line_pt=1.0)
    add_textbox(slide, QX, Y + TOTAL_H/2 - 0.35, QW, 0.7, [
        {"text": "사용자\n질문", "sz": 10, "bold": True, "color_hex": NAVY,
         "align": PP_ALIGN.CENTER}
    ])

    add_textbox(slide, QX + QW + 0.05, Y + TOTAL_H/2 - 0.22, 0.4, 0.44, [
        {"text": "→", "sz": 12, "bold": True, "color_hex": GRAY20,
         "align": PP_ALIGN.CENTER}
    ])

    # 사용자 질문 分析 박스 (하나로 통합)
    LX = QX + QW + 0.5
    LW = 3.2
    add_rect(slide, LX, Y, LW, TOTAL_H, fill_hex=LBLUE, line_hex=BLUE, line_pt=0.8)
    add_textbox(slide, LX, Y + TOTAL_H/2 - 0.55, LW, 1.1, [
        {"text": "사용자 질문 分析", "sz": 9.5, "bold": True, "color_hex": NAVY,
         "align": PP_ALIGN.CENTER},
        {"text": "(GPT-OSS-20B)", "sz": 8, "color_hex": TEAL,
         "align": PP_ALIGN.CENTER},
        {"text": "의도+엔티티", "sz": 7.5, "color_hex": TEAL,
         "align": PP_ALIGN.CENTER},
    ])

    # LLM → 의도 화살표 (rows 0-2)
    LAX = LX + LW + 0.05
    for r in range(3):
        add_textbox(slide, LAX, row_y(r) + H_ROW/2 - 0.22, 0.4, 0.44, [
            {"text": "→", "sz": 11, "bold": True, "color_hex": GRAY20,
             "align": PP_ALIGN.CENTER}
        ])

    # 의도 박스 (risk rows 0-1, law row 2, design row 3)
    IW = 2.0
    IX = LAX + 0.45
    risk_h = 2 * H_ROW + V_GAP
    add_rect(slide, IX, row_y(0), IW, risk_h,
             fill_hex=TABLE_COLORS["accidents"][1], line_hex=TABLE_COLORS["accidents"][0], line_pt=0.8)
    add_textbox(slide, IX, row_y(0) + risk_h/2 - 0.22, IW, 0.44, [
        {"text": "risk", "sz": 10, "bold": True,
         "color_hex": TABLE_COLORS["accidents"][0], "align": PP_ALIGN.CENTER}
    ])
    add_rect(slide, IX, row_y(2), IW, H_ROW,
             fill_hex=TABLE_COLORS["laws"][1], line_hex=TABLE_COLORS["laws"][0], line_pt=0.8)
    add_textbox(slide, IX, row_y(2) + H_ROW/2 - 0.22, IW, 0.44, [
        {"text": "law", "sz": 10, "bold": True,
         "color_hex": TABLE_COLORS["laws"][0], "align": PP_ALIGN.CENTER}
    ])
    add_rect(slide, IX, row_y(3), IW, H_ROW,
             fill_hex=TABLE_COLORS["designs"][1], line_hex=TABLE_COLORS["designs"][0], line_pt=0.8)
    add_textbox(slide, IX, row_y(3) + H_ROW/2 - 0.22, IW, 0.44, [
        {"text": "design", "sz": 10, "bold": True,
         "color_hex": TABLE_COLORS["designs"][0], "align": PP_ALIGN.CENTER}
    ])

    # 의도/엔티티 → DB 화살표
    AX = IX + IW + 0.05
    AW_ARR = 0.4
    for r in range(4):
        add_textbox(slide, AX, row_y(r) + H_ROW/2 - 0.22, AW_ARR, 0.44, [
            {"text": "→", "sz": 11, "bold": True, "color_hex": GRAY20,
             "align": PP_ALIGN.CENTER}
        ])

    # DB 박스
    DBX = AX + AW_ARR + 0.05
    DBW = 4.2
    for name, r in [("accidents", 0), ("chemicals", 1), ("laws", 2), ("designs", 3)]:
        dc, lc = TABLE_COLORS[name]
        add_rect(slide, DBX, row_y(r), DBW, H_ROW, fill_hex=lc, line_hex=dc, line_pt=0.8)
        add_textbox(slide, DBX + 0.1, row_y(r) + H_ROW/2 - 0.22, DBW - 0.2, 0.44, [
            {"text": name + " DB", "sz": 9, "bold": True, "color_hex": dc,
             "align": PP_ALIGN.CENTER}
        ])

    # 수직 연결선 + 화살표
    VX = DBX + DBW + 0.05
    add_rect(slide, VX, row_y(0), 0.08, TOTAL_H, fill_hex=GRAY50)
    add_textbox(slide, VX + 0.08, Y + TOTAL_H/2 - 0.22, 0.45, 0.44, [
        {"text": "→", "sz": 13, "bold": True, "color_hex": GRAY20,
         "align": PP_ALIGN.CENTER}
    ])

    # Hybrid Search 박스
    HSX = VX + 0.08 + 0.5
    HSW = 5.2
    add_rect(slide, HSX, Y, HSW, TOTAL_H, fill_hex=LBLUE, line_hex=BLUE, line_pt=1.0)
    add_textbox(slide, HSX, Y + TOTAL_H/2 - 0.45, HSW, 0.9, [
        {"text": "① Hybrid Search", "sz": 11, "bold": True, "color_hex": NAVY,
         "align": PP_ALIGN.CENTER},
        {"text": "BGE-M3 + BM25 (0.5 / 0.5)", "sz": 8.5, "color_hex": TEAL,
         "align": PP_ALIGN.CENTER},
    ])

    # Re-ranking 박스
    A2X = HSX + HSW + 0.05
    add_textbox(slide, A2X, Y + TOTAL_H/2 - 0.22, 0.45, 0.44, [
        {"text": "→", "sz": 13, "bold": True, "color_hex": GRAY20,
         "align": PP_ALIGN.CENTER}
    ])
    RRX = A2X + 0.5
    RRW = 5.0
    add_rect(slide, RRX, Y, RRW, TOTAL_H, fill_hex=LGREEN, line_hex=GREEN, line_pt=1.0)
    add_textbox(slide, RRX, Y + TOTAL_H/2 - 0.45, RRW, 0.9, [
        {"text": "② Re-ranking", "sz": 11, "bold": True, "color_hex": GREEN,
         "align": PP_ALIGN.CENTER},
        {"text": "BGE-Reranker-v2-m3", "sz": 8.5, "color_hex": GREEN,
         "align": PP_ALIGN.CENTER},
    ])

    # 최종 반환 박스
    A3X = RRX + RRW + 0.05
    add_textbox(slide, A3X, Y + TOTAL_H/2 - 0.22, 0.45, 0.44, [
        {"text": "→", "sz": 13, "bold": True, "color_hex": GRAY20,
         "align": PP_ALIGN.CENTER}
    ])
    FX = A3X + 0.5
    FW = 33.0 - FX
    add_rect(slide, FX, Y, FW, TOTAL_H, fill_hex="FEF9E7", line_hex=ORANGE, line_pt=1.0)
    add_textbox(slide, FX, Y + TOTAL_H/2 - 0.55, FW, 1.1, [
        {"text": "최종 반환", "sz": 10, "bold": True, "color_hex": ORANGE,
         "align": PP_ALIGN.CENTER},
        {"text": "관련 문서 청크", "sz": 8.5, "color_hex": GRAY20,
         "align": PP_ALIGN.CENTER},
        {"text": "예: accidents 6건\nlaws 3건", "sz": 7.5, "color_hex": GRAY50,
         "align": PP_ALIGN.CENTER},
    ])

    # LLM 엔티티 결과 노트
    ent_y = Y + TOTAL_H + 0.25
    add_rect(slide, 0.8, ent_y, 32.2, 0.7, fill_hex="F8F9FA", line_hex=GRAY50, line_pt=0.5)
    add_textbox(slide, 1.0, ent_y + 0.13, 32.0, 0.5, [
        {"text": "LLM 분석 출력 ②: target_material / target_equipment  →  Hybrid Search 시 해당 필드 벡터 정밀 매칭에 활용",
         "sz": 9, "color_hex": GRAY20}
    ])

    Y = ent_y + 0.85

    # === 단계별 역할 요약 ===
    section_header(slide, 0.8, Y, 32, "단계별 역할 요약")
    Y += 0.65

    summary_data = [
        ["단계",             "역할",                                     "출력"],
        ["LLM 분석",         "의도 판단 + 물질·설비명 추출",              "intents / target_material / target_equipment"],
        ["① Hybrid Search",  "의미(BGE-M3) + 키워드(BM25) 혼합 검색",    "테이블별 상위 N건 후보"],
        ["② Re-ranking",     "(질문, 문서) 쌍 정밀 분석 → 최종 선별",    "임계치 통과 문서만 반환"],
    ]
    add_table(slide, summary_data,
              l=0.8, t=Y, w=32.0, h=2.5,
              col_widths=[4, 7, 8],
              font_sz=8.5, header_sz=9.5,
              h_data=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# Slide 5: 사용자 질문 分析
# ════════════════════════════════════════════════════════════════════════════

def build_slide16(slide):
    clear_slide(slide)
    update_title(slide, "RAG 구축 - 사용자 질문 分析")
    update_titlebar(slide,
        "사용자 질문 分析",
        "LLM이 질문에서 검색에 필요한 정보를 한 번에 추출 — 룰기반 대비 장점")

    Y = 6.3

    # === 왜 LLM으로 질문을 분류하는가? ===
    section_header(slide, 0.8, Y, 32, "왜 LLM으로 질문을 분류하는가?")
    Y += 0.65

    add_textbox(slide, 0.8, Y, 32.0, 0.35, [
        {"text": "해야 할 작업: 사용자 질문에서 의도(어느 DB를 검색할지)와 핵심 정보(물질명·설비명)를 동시에 파악",
         "sz": 9, "bold": True, "color_hex": NAVY}
    ])
    Y += 0.45

    add_textbox(slide, 0.8, Y, 32.0, 1.6, [
        {"text": "룰기반 방식(키워드/패턴 매칭)으로 하면?",
         "sz": 9, "bold": True, "color_hex": TEAL},
        {"text": "  - 방법: \"폭발\"·\"위험\" 등 키워드 존재 시 risk로 판단, 사전 정의된 물질명 목록과 패턴 매칭으로 추출",
         "sz": 8.5, "color_hex": GRAY20},
        {"text": "  - 한계 ①: 화학공정 물질·설비명은 수천 종 이상 → 사전 완성 현실적으로 불가, 신조어·약어 대응 불가",
         "sz": 8.5, "color_hex": GRAY20},
        {"text": "  - 한계 ②: 예외 케이스 발생마다 규칙 추가 → 코드 복잡도 무한 증가, 유지보수 불가",
         "sz": 8.5, "color_hex": GRAY20},
    ])
    Y += 1.7

    add_textbox(slide, 0.8, Y, 32.0, 0.8, [
        {"text": "LLM + Few-shot으로 해결:",
         "sz": 9, "bold": True, "color_hex": GREEN},
        {"text": "  예시 Q&A 몇 개만 프롬프트에 포함 → 새로운 표현·복합 질문·화학 전문 용어도 문맥 파악하여 처리 / 별도 학습 불필요",
         "sz": 8.5, "color_hex": GRAY20},
    ])
    Y += 1.0

    # === LLM 호출 구조 및 예시 ===
    section_header(slide, 0.8, Y, 32, "LLM 호출 구조 및 예시")
    Y += 0.65

    # 왼쪽: 프롬프트 구조 코드박스
    LW = 18.0
    RW = 32.0 - LW - 0.6
    RX = 0.8 + LW + 0.6

    add_code_box(slide, 0.8, Y, LW, 5.8, [
        "[시스템 프롬프트]",
        "역할: 너는 화학공정 안전 전문 AI야.",
        "      사용자 질문에서 의도와 핵심 정보를 추출해.",
        '{\"intents\": ..., \"target_material\": \"...\", \"target_equipment\": \"...\"} JSON만 반환',
        "",
        "[Few-shot 예시]",
        '  Q: "안전밸브 미설치 시 과태료가 얼마인지 알려줘."',
        '  A: {\"intents\": \"law\", \"target_material\": \"None\", \"target_equipment\": \"\uc548\uc804\ubc38\ube0c\"}',
        "",
        '  Q: "V-101 탱크 내부 압력 급상승 시 폭발 가능성?"',
        '  A: {\"intents\": \"risk\", \"target_material\": \"None\", \"target_equipment\": \"V-101 탱크\"}',
        "",
        '  Q: "이 설비 위험해?"',
        '  A: {\"intents\": \"risk\", \"target_material\": \"None\", \"target_equipment\": \"None\"}',
        "",
        "[유저 프롬프트]",
        "  톨루엔을 취급하는 반응기의 폭발 위험성을 알려줘.",
    ], font_sz=8.5)

    # 오른쪽: LLM 출력 + 활용 표
    add_textbox(slide, RX, Y, RW, 0.4, [
        {"text": "LLM 출력", "sz": 9, "bold": True, "color_hex": ORANGE}
    ])
    add_code_box(slide, RX, Y + 0.45, RW, 1.4, [
        '{',
        '  \"intents\": \"risk\",',
        '  \"target_material\": "톨루엔",',
        '  \"target_equipment\": "반응기"',
        '}',
    ], font_sz=9)

    usage_data = [
        ["출력 필드",         "활용"],
        ["intents",          "검색 대상 DB 결정\n\"risk\" → accidents + chemicals"],
        ["target_material",  "accidents DB material 벡터 정밀 매칭"],
        ["target_equipment", "accidents DB equipment 벡터 정밀 매칭"],
        ["값이 None이면",     "해당 필드 매칭 생략\n→ 전체 text_vector로만 검색"],
    ]
    add_table(slide, usage_data,
              l=RX, t=Y + 2.05, w=RW, h=3.6,
              col_widths=[4, 7.5],
              font_sz=8.5, header_sz=9.5,
              h_data=PP_ALIGN.LEFT)
# ════════════════════════════════════════════════════════════════════════════
# Slide 6: 검색 — Semantic Search
# ════════════════════════════════════════════════════════════════════════════

def build_slide_semantic(slide):
    clear_slide(slide)
    update_title(slide, "RAG 구축 - 검색 (Semantic Search)")
    update_titlebar(slide,
        "검색 (1단계)  —  Semantic Search",
        "텍스트를 벡터로 변환해 의미 기반으로 유사 문서를 찾는다")

    Y = 6.3

    # === Semantic Search란? ===
    section_header(slide, 0.8, Y, 32, "Semantic Search란?")
    Y += 0.65

    add_textbox(slide, 0.8, Y, 32.0, 0.55, [
        {"text": "텍스트를 고차원 숫자 배열(벡터)로 변환해 의미가 비슷한 문서를 찾는 검색 방식.",
         "sz": 9.5, "color_hex": GRAY20},
        {"text": "핵심 아이디어: \"비슷한 의미 = 벡터 공간에서 가까운 위치\"",
         "sz": 9.5, "bold": True, "color_hex": NAVY},
    ])
    Y += 0.75

    add_code_box(slide, 0.8, Y, 32.0, 2.1, [
        '예시:',
        '  "톨루엔 반응기 폭발 위험"   →  [0.12, -0.85,  0.43, ...]  (1,024차원)',
        '  "인화성 물질 용기 화재 사고" →  [0.10, -0.81,  0.39, ...]  ← 의미 유사 → 벡터 가까움',
        '  "펌프 유량 측정 방법"        →  [0.91,  0.23, -0.67, ...]  ← 의미 무관 → 벡터 멂',
        '',
        '→ 표현이 달라도 의미가 유사하면 검색 가능 (키워드 완전 일치 불필요)',
    ], font_sz=8.5)
    Y += 2.3

    # === 메커니즘 단계 (박스+화살표 도식) ===
    section_header(slide, 0.8, Y, 32, "메커니즘 단계")
    Y += 0.65

    steps_sem = [
        ("질문 텍스트", PGRAY, PGRAY2, GRAY20),
        ("BGE-M3\n임베딩 모델", LBLUE, BLUE, NAVY),
        ("1,024차원\n벡터 변환", LBLUE, BLUE, NAVY),
        ("코사인 유사도\n계산", LBLUE, BLUE, NAVY),
        ("상위 N건\n후보 추출", LGREEN, GREEN, GREEN),
    ]
    BW_S = (32.0 - 4 * 0.5) / 5  # ≈ 6.0
    AW_S = 0.5
    BH_S = 1.7
    SX = 0.8
    for i, (label, fill, border, tc) in enumerate(steps_sem):
        add_rect(slide, SX, Y, BW_S, BH_S, fill_hex=fill, line_hex=border, line_pt=1.0)
        add_textbox(slide, SX + 0.1, Y + 0.2, BW_S - 0.2, BH_S - 0.2, [
            {"text": label, "sz": 9, "bold": True, "color_hex": tc, "align": PP_ALIGN.CENTER}
        ])
        SX += BW_S
        if i < 4:
            add_textbox(slide, SX, Y + BH_S / 2 - 0.22, AW_S, 0.44, [
                {"text": "→", "sz": 12, "bold": True, "color_hex": GRAY20, "align": PP_ALIGN.CENTER}
            ])
            SX += AW_S

    Y += BH_S + 0.2

    add_textbox(slide, 0.8, Y, 32.0, 0.7, [
        {"text": "코사인 유사도 = cos(θ) = (A·B) / (|A||B|)   — 1에 가까울수록 의미 유사, 0이면 무관",
         "sz": 8.5, "color_hex": GRAY50},
        {"text": "※ 이 도식은 Semantic Search 단계만 표현. BM25 가중치 합산은 다음 슬라이드(BM25 + Hybrid)에서 다룸.",
         "sz": 8, "color_hex": GRAY50},
    ])
    Y += 0.9
    section_header(slide, 0.8, Y, 32, "BGE-M3 선택 이유 + 강점·약점")
    Y += 0.65

    HALF_W_S = 15.8
    GAP_S = 0.6
    X2_S = 0.8 + HALF_W_S + GAP_S

    # 왼쪽: 선택 이유
    add_textbox(slide, 0.8, Y, HALF_W_S, 1.5, [
        {"text": "BGE-M3 선택 이유", "sz": 9, "bold": True, "color_hex": TEAL},
        {"text": "▸  한국어·영어 동시 지원 → 법령명(한국어) + 화학물질명(영문 혼재) 모두 처리",
         "sz": 8.5, "color_hex": GRAY20},
        {"text": "▸  1,024차원: 화학공정 전문 문서의 복잡한 의미 표현에 충분한 표현력",
         "sz": 8.5, "color_hex": GRAY20},
    ])

    # 오른쪽: 강점·약점
    sw_data = [
        ["",       "내용"],
        ["강점",   "\"폭발 사고\"와 \"화재 사고\"처럼 표현이 다른 유사 사고 검색 가능"],
        ["약점",   "\"산안법 제36조\", \"V-101\" 같은 고유 번호·코드는 의미 벡터로 구분 어려움"],
    ]
    add_table(slide, sw_data,
              l=X2_S, t=Y, w=HALF_W_S, h=1.8,
              col_widths=[1.8, 10],
              font_sz=8.5, header_sz=9,
              h_data=PP_ALIGN.LEFT)

    Y += 2.0
    add_textbox(slide, 0.8, Y, 32.0, 0.4, [
        {"text": "→ 이 약점을 보완하기 위해 BM25를 함께 사용 (다음 슬라이드)",
         "sz": 9, "bold": True, "color_hex": NAVY}
    ])


# ════════════════════════════════════════════════════════════════════════════
# Slide 7: 검색 — BM25 + Hybrid Score
# ════════════════════════════════════════════════════════════════════════════

def build_slide_bm25(slide):
    clear_slide(slide)
    update_title(slide, "RAG 구축 - 검색 (BM25 + Hybrid)")
    update_titlebar(slide,
        "검색 (2단계)  —  BM25 + Hybrid Score",
        "키워드 빈도 기반 BM25와 Semantic을 합산해 상호 약점 보완")

    Y = 6.3

    # === BM25란? ===
    section_header(slide, 0.8, Y, 32, "BM25란?")
    Y += 0.65

    add_textbox(slide, 0.8, Y, 32.0, 0.4, [
        {"text": "단어의 출현 빈도를 기반으로 문서의 관련도를 계산하는 고전 정보 검색 알고리즘. 두 가지 핵심 개념:",
         "sz": 9, "color_hex": GRAY20}
    ])
    Y += 0.55

    tf_idf_data = [
        ["개념",  "의미",                                  "예시"],
        ["TF\n(Term Frequency)",
         "이 문서에 해당 단어가 얼마나 자주 나오는가",
         "\"톨루엔\"이 5번 등장 → TF 높음"],
        ["IDF\n(Inverse Document Frequency)",
         "전체 문서 중 이 단어가 얼마나 희귀한가",
         "\"을/이\"는 흔함 → IDF 낮음\n\"톨루엔\"은 희귀 → IDF 높음"],
    ]
    add_table(slide, tf_idf_data,
              l=0.8, t=Y, w=32.0, h=2.5,
              col_widths=[3.5, 8, 8],
              font_sz=8.5, header_sz=9.5,
              h_data=PP_ALIGN.LEFT)
    Y += 2.65

    add_code_box(slide, 0.8, Y, 32.0, 0.85, [
        "BM25 점수 = TF(단어, 문서) × IDF(단어)  →  0~1로 정규화",
        "→ 흔한 조사(\"을\", \"이\")는 점수 낮음 / 전문 용어(\"산안법 제36조\", \"V-101\")는 점수 높음",
    ], font_sz=9)
    Y += 1.05

    # === BM25 vs Semantic 비교 + Hybrid 선택 이유 ===
    section_header(slide, 0.8, Y, 32, "BM25 vs Semantic 비교 → 왜 Hybrid?")
    Y += 0.65

    cmp_data = [
        ["",        "Semantic Search",                       "BM25"],
        ["강점",    "표현이 달라도 의미 유사 문서 탐색",      "전문 용어·법조항 번호 정확 매칭"],
        ["약점",   "고유 번호·코드명 구분 어려움",           "표현이 다른 유사 사고 탐색 불가"],
        ["결론",    "→ 두 방식이 서로의 약점을 정확히 보완  →  Hybrid로 결합", ""],
    ]
    add_table(slide, cmp_data,
              l=0.8, t=Y, w=32.0, h=2.5,
              col_widths=[2.5, 9, 9],
              font_sz=8.5, header_sz=9.5,
              h_data=PP_ALIGN.LEFT)
    Y += 2.7

    add_textbox(slide, 0.8, Y, 32.0, 0.4, [
        {"text": "▸  BM25 선택 추가 이유: LanceDB에 기본 내장 → 별도 구현 불필요",
         "sz": 8.5, "color_hex": GRAY50}
    ])
    Y += 0.6

    # === Hybrid Score 합산 및 후보 선별 ===
    section_header(slide, 0.8, Y, 32, "Hybrid Score 합산 및 후보 선별")
    Y += 0.65

    add_code_box(slide, 0.8, Y, 32.0, 0.7, [
        "Hybrid Score = 0.5 × Semantic 점수 + 0.5 × BM25 점수  →  상위 N건 후보 선별",
    ], font_sz=9.5)
    Y += 0.9

    cand_data = [
        ["테이블",    "1차 상위 후보 수", "비고"],
        ["accidents", "300건",           "re-ranking에서 정밀 선별할 것이므로 넉넉하게"],
        ["chemicals", "3건",             "물질 1종 = 1청크 구조 → 소수로 충분"],
        ["laws",      "10건",            ""],
        ["designs",   "10건",            ""],
    ]
    add_table(slide, cand_data,
              l=0.8, t=Y, w=32.0, h=2.5,
              col_widths=[3, 4, 12],
              font_sz=8.5, header_sz=9.5,
              h_data=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# Slide 8: 검색 — 2차 Cross-Encoder Re-ranking
# ════════════════════════════════════════════════════════════════════════════

def build_slide17(slide):
    clear_slide(slide)
    update_title(slide, "RAG 구축 - 검색 (Re-ranking)")
    update_titlebar(slide,
        "검색 (3단계)  —  Cross-Encoder Re-ranking",
        "1차 대량 후보를 (질문, 문서) 쌍으로 정밀 재정렬 → 최종 소수 문서 선별")

    Y = 6.3

    # === 왜 Re-ranking이 필요한가? ===
    section_header(slide, 0.8, Y, 32, "왜 Re-ranking이 필요한가?")
    Y += 0.65

    add_textbox(slide, 0.8, Y, 32.0, 2.2, [
        {"text": "1차 Hybrid Search의 구조적 한계:",
         "sz": 9, "bold": True, "color_hex": TEAL},
        {"text": "  ▸  Semantic Search: 질문 벡터와 문서 벡터를 독립적으로 인코딩한 뒤 유사도 비교 → \"이 질문에 대해 이 문서가 적합한가\"를 모델이 직접 판단하지 못함",
         "sz": 8.5, "color_hex": GRAY20},
        {"text": "  ▸  BM25: 단어 빈도만 보므로 \"반응기\"라는 단어가 많은 무관한 문서가 상위에 올 수 있음",
         "sz": 8.5, "color_hex": GRAY20},
        {"text": "  → 결과적으로 1차에서는 관련 있어 보이지만 실제로는 무관한 문서가 섞여 들어옴",
         "sz": 8.5, "bold": True, "color_hex": ORANGE},
        {"text": "Cross-Encoder Re-ranking: (질문, 문서) 쌍을 함께 모델에 입력 → 질문 맥락에서 이 문서가 실제로 관련 있는지 직접 판단",
         "sz": 8.5, "bold": True, "color_hex": GREEN},
    ])
    Y += 2.4

    # === Cross-Encoder 메커니즘 ===
    section_header(slide, 0.8, Y, 32, "Cross-Encoder 메커니즘")
    Y += 0.65

    # Cross-Encoder 메커니즘: 박스+화살표 도식
    steps_ce = [
        ("Hybrid Search\n후보 N건", PGRAY, PGRAY2, GRAY20),
        ("(질문, 문서)\n쌍 생성", LBLUE, BLUE, NAVY),
        ("BGE-Reranker\n-v2-m3", LBLUE, BLUE, NAVY),
        ("관련도 점수\n(0~1) 계산", LBLUE, BLUE, NAVY),
        ("임계치 통과\n문서 선별", LGREEN, GREEN, GREEN),
    ]
    BW_CE = (32.0 - 4 * 0.5) / 5
    AW_CE = 0.5
    BH_CE = 1.7
    CX = 0.8
    for i, (label, fill, border, tc) in enumerate(steps_ce):
        add_rect(slide, CX, Y, BW_CE, BH_CE, fill_hex=fill, line_hex=border, line_pt=1.0)
        add_textbox(slide, CX + 0.1, Y + 0.2, BW_CE - 0.2, BH_CE - 0.2, [
            {"text": label, "sz": 9, "bold": True, "color_hex": tc, "align": PP_ALIGN.CENTER}
        ])
        CX += BW_CE
        if i < 4:
            add_textbox(slide, CX, Y + BH_CE / 2 - 0.22, AW_CE, 0.44, [
                {"text": "→", "sz": 12, "bold": True, "color_hex": GRAY20, "align": PP_ALIGN.CENTER}
            ])
            CX += AW_CE
    Y += BH_CE + 0.2

    add_textbox(slide, 0.8, Y, 32.0, 0.55, [
        {"text": "쌍 생성: (질문, 문서1), (질문, 문서2), ..., (질문, 문서N)",
         "sz": 8.5, "color_hex": GRAY50},
        {"text": "Cross-Encoder: 질문+문서를 함께 입력 → Semantic Search와 달리 두 텍스트 간 상호작용 직접 반영",
         "sz": 8.5, "color_hex": GRAY50},
    ])
    Y += 0.85

    # === 테이블별 최종 반환 ===
    section_header(slide, 0.8, Y, 32, "테이블별 최종 반환")
    Y += 0.65

    limit_data = [
        ["테이블",    "1차 후보", "임계치", "최종 반환", "임계치 설정 이유"],
        ["accidents", "300건",   "0.7",   "최대 6건",  "사고 정보는 정확도가 중요 → 높은 임계치"],
        ["chemicals", "3건",     "0.7",   "최대 1건",  "사고 정보는 정확도가 중요 → 높은 임계치"],
        ["laws",      "10건",    "0.5",   "최대 3건",  "법령 문서는 수가 적어 넓게 수집"],
        ["designs",   "10건",    "0.5",   "최대 3건",  "설계 기준은 수가 적어 넓게 수집"],
    ]
    add_table(slide, limit_data,
              l=0.8, t=Y, w=32.0, h=3.0,
              col_widths=[3, 3, 2.5, 3, 9],
              font_sz=8.5, header_sz=9.5,
              h_data=PP_ALIGN.LEFT,
              h_first=PP_ALIGN.CENTER)


def build_slide18(slide):
    clear_slide(slide)
    update_title(slide, "RAG 구축 - 동작 예시")
    update_titlebar(slide,
        "동작 예시  —  실제 처리 흐름",
        "예시 질문이 LLM 분석 → Hybrid Search → Re-ranking을 거치는 과정 추적")

    Y = 6.3

    # ── 예시 질문 ──
    q_box = add_rect(slide, 0.8, Y, 32.0, 0.85,
                     fill_hex="EBF3FB", line_hex=BLUE, line_pt=1)
    tf = q_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = '예시 질문:  "톨루엔을 취급하는 반응기의 폭발 위험성을 알려줘."'
    r.font.size  = Pt(11)
    r.font.bold  = True
    r.font.color.rgb = RGBColor.from_string(NAVY)
    Y += 1.0

    # ── 3단계 가로 박스 ──
    section_header(slide, 0.8, Y, 32, "단계별 처리")
    Y += 0.65

    BW = (32.0 - 0.3 * 2) / 3   # ≈ 10.47
    BH = 4.4
    BXSTART = 0.8
    BOX_Y = Y

    step_boxes = [
        {
            "title":   "① LLM 분석",
            "fill":    LBLUE, "border": BLUE, "title_color": NAVY,
            "in_lines":  ['"톨루엔을 취급하는 반응기의', '폭발 위험성을 알려줘."'],
            "out_lines": ['{"intents": ["risk"],',
                          ' "target_material": "톨루엔",',
                          ' "target_equipment": "반응기"}',
                          '→ accidents + chemicals 검색'],
        },
        {
            "title":   "② Hybrid Search",
            "fill":    "EBF3FB", "border": TEAL, "title_color": NAVY,
            "in_lines":  ['질문 벡터 (1,024차원)', '+ BM25 키워드 점수'],
            "out_lines": ['accidents 상위 300건 후보',
                          '(BGE-M3 × 0.5 + BM25 × 0.5)'],
        },
        {
            "title":   "③ Cross-Encoder Re-ranking",
            "fill":    LGREEN, "border": GREEN, "title_color": GREEN,
            "in_lines":  ['(질문, 문서) 300쌍', 'BGE-Reranker-v2-m3'],
            "out_lines": ['score ≥ 0.7 통과', '→ 최종 6건 반환'],
        },
    ]

    for i, box in enumerate(step_boxes):
        BX = BXSTART + i * (BW + 0.3)
        add_rect(slide, BX, BOX_Y, BW, BH,
                 fill_hex=box["fill"], line_hex=box["border"], line_pt=1.0)
        # 제목
        add_textbox(slide, BX + 0.15, BOX_Y + 0.15, BW - 0.3, 0.55, [
            {"text": box["title"], "sz": 10, "bold": True,
             "color_hex": box["title_color"], "align": PP_ALIGN.CENTER}
        ])
        # 구분선
        add_rect(slide, BX + 0.2, BOX_Y + 0.75, BW - 0.4, 0.04, fill_hex=box["border"])
        # 입력
        add_textbox(slide, BX + 0.2, BOX_Y + 0.85, BW - 0.4, 0.35, [
            {"text": "▸ 입력", "sz": 8.5, "bold": True, "color_hex": NAVY}
        ])
        add_code_box(slide, BX + 0.2, BOX_Y + 1.2, BW - 0.4, 0.85,
                     box["in_lines"], font_sz=8.5)
        # 출력
        add_textbox(slide, BX + 0.2, BOX_Y + 2.15, BW - 0.4, 0.35, [
            {"text": "▸ 출력", "sz": 8.5, "bold": True, "color_hex": GREEN}
        ])
        add_code_box(slide, BX + 0.2, BOX_Y + 2.5, BW - 0.4, 1.75,
                     box["out_lines"], font_sz=8.5)

    # 박스 사이 화살표
    for i in range(2):
        ARX = BXSTART + (i + 1) * BW + i * 0.3
        add_textbox(slide, ARX, BOX_Y + BH / 2 - 0.22, 0.3, 0.44, [
            {"text": "→", "sz": 12, "bold": True,
             "color_hex": GRAY20, "align": PP_ALIGN.CENTER}
        ])

    Y = BOX_Y + BH + 0.35

    # ── 최종 결과 카드 ──
    section_header(slide, 0.8, Y, 32, "최종 반환 결과  (score ≥ 0.7 통과, accidents 상위 3건)")
    Y += 0.65

    results = [
        ("1위", "0.961",
         "회분식 반응기에 톨루엔 투입 후 맨홀 덮개를 닫는 순간\n내부 증기 점화 → 폭발,  사망 1명",
         "폭발", "인화성물질 취급방법 불량"),
        ("2위", "0.934",
         "반응기 내 톨루엔 유증기 폭발\n인근 작업자 2명 화상 (전치 3개월 이상)",
         "폭발", "원인 미상 점화원"),
        ("3위", "0.912",
         "알키드수지·톨루엔 혼합물 반응기 누출 → 폭발·화재\n사망 1명",
         "화재", "비점 이상 톨루엔 주입"),
    ]

    CW = BW
    for idx, (rank, score, content, atype, cause) in enumerate(results):
        CX = 0.8 + idx * (CW + 0.3)
        add_rect(slide, CX, Y, CW, 0.5, fill_hex=DARK)
        add_textbox(slide, CX + 0.1, Y + 0.07, CW - 0.2, 0.4, [
            {"text": f"{rank}  (Re-ranking score: {score})",
             "sz": 9, "bold": True,
             "color_hex": ORANGE if idx == 0 else LGRAY}
        ])
        add_rect(slide, CX, Y + 0.5, CW, 3.5,
                 fill_hex="F7F7F7", line_hex=PGRAY2, line_pt=0.5)
        add_textbox(slide, CX + 0.15, Y + 0.65, CW - 0.3, 3.2, [
            {"text": f"[사고내용]  {content}", "sz": 8.5, "color_hex": GRAY20},
            {"text": f"[관련물질] 톨루엔   [관련설비] 반응기", "sz": 8, "color_hex": "444444"},
            {"text": f"[사고유형] {atype}   [사고원인] {cause}", "sz": 8, "color_hex": "444444"},
        ])

    Y += 4.1
    add_textbox(slide, 0.8, Y, 32.0, 0.4, [
        {"text": "* score는 Re-ranking 점수. 실제 실행 시 달라질 수 있음.",
         "sz": 8, "color_hex": GRAY50}
    ])


# ════════════════════════════════════════════════════════════════════════════
# Main
# ════════════════════════════════════════════════════════════════════════════

# ════════════════════════════════════════════════════════════════════════════
# Slide 1: 수집 데이터 목록 업데이트 (사용자 추가 슬라이드)
# ════════════════════════════════════════════════════════════════════════════

def update_slide1_table(slide):
    """Slide 1의 기존 표를 실제 corpus 5종 데이터로 교체"""
    for shape in list(slide.shapes):
        if shape.shape_type == 19:  # TABLE
            shape._element.getparent().remove(shape._element)
            break

    new_data = [
        ["No.", "데이터 유형", "핵심 지식 내용", "주요 활용 목적", "데이터 건수"],
        ["1", "사고사례\n(accidents)",
         "국내외 화학공정 사고 이력\n사고내용, 관련물질·설비, 사고유형, 원인",
         "공정 위험성 질문에 대한\n유사 사고 사례 검색·제시",
         "1,348건\n(xlsx 3종)"],
        ["2", "화학물질정보\n(chemicals)",
         "460종 화학물질의 MSDS 정보\n위험성, 물성, 취급주의사항",
         "특정 물질의 위험성·취급 정보 제공",
         "501종\n(pdf 1종)"],
        ["3", "법령\n(laws)",
         "산업안전보건법, 고압가스안전관리법 등\n국내 법령·고시·기준",
         "법규 위반 여부 판단 및\n해당 조항 근거 제시",
         "2,099페이지\n(pdf 43종)"],
        ["4", "설계지침\n(designs)",
         "KOSHA Guide, KGS 코드 등\n설계·운전 기술지침",
         "설계 오류 여부 판단 및\n개선 기준 제시",
         "1,983페이지\n(pdf 61종)"],
        ["5", "화공 기초지식\n(basics)",
         "화공 전공서, 논문, 보고서\n화학공정 원리·이론",
         "배경 지식 기반\n공정 위험 설명 보완",
         "8,165페이지\n(pdf 15종)"],
    ]
    add_table(slide, new_data,
              l=0.44, t=3.90, w=32.99, h=14.13,
              col_widths=[1.2, 3.5, 9.0, 7.0, 3.5],
              font_sz=9.5, header_sz=10.5,
              h_data=PP_ALIGN.LEFT,
              h_first=PP_ALIGN.CENTER,
              h_header=PP_ALIGN.CENTER)
    print("  Slide 1 표 교체 완료")


def main():
    prs = Presentation(OUTPUT if os.path.isfile(OUTPUT) else INPUT)
    n = len(prs.slides)
    print(f"입력 파일: {OUTPUT if os.path.isfile(OUTPUT) else INPUT}  ({n} 슬라이드)")

    if n >= 7:
        # 목표: 11슬라이드 (사용자 2 + RAG 9)
        while len(prs.slides) < 11:
            duplicate_slide(prs, 2)

        print("Slide 1 (수집 데이터 목록) 표 업데이트...")
        update_slide1_table(prs.slides[0])

        print("RAG Slide 1 (DB 구축 흔름) 구축...")
        build_slide14(prs.slides[2])

        print("RAG Slide 2 (DB 구축 상세) 구축...")
        build_slide14b(prs.slides[3])

        print("RAG Slide 3 (DB 구축 현황) 구축...")
        build_slide15(prs.slides[4])

        print("RAG Slide 4 (검색 파이프라인) 구축...")
        build_slide_pipeline(prs.slides[5])

        print("RAG Slide 5 (사용자 질문 分析) 구축...")
        build_slide16(prs.slides[6])

        print("RAG Slide 6 (Semantic Search) 구축...")
        build_slide_semantic(prs.slides[7])

        print("RAG Slide 7 (BM25 + Hybrid) 구축...")
        build_slide_bm25(prs.slides[8])

        print("RAG Slide 8 (Re-ranking) 구축...")
        build_slide17(prs.slides[9])

        print("RAG Slide 9 (동작 예시) 구축...")
        build_slide18(prs.slides[10])

    else:
        # 원본 "RAG 부분.pptx" (2슬라이드) 기반 플로우
        print("Slide 14 (Vector DB 구성) 구축...")
        build_slide14(prs.slides[0])

        print("Slide 14b (DB 구축 상세) 구축...")
        s14b = duplicate_slide(prs, 0)
        build_slide14b(s14b)

        print("Slide 15 (DB 구축 현황) 구축...")
        s15 = duplicate_slide(prs, 0)
        build_slide15(s15)

        print("Slide pipeline (검색 파이프라인) 구축...")
        sp = duplicate_slide(prs, 0)
        build_slide_pipeline(sp)

        print("Slide 16 (질문 분석) 구축...")
        s16 = duplicate_slide(prs, 0)
        build_slide16(s16)

        print("Slide 17 (검색 과정) 구축...")
        s17 = duplicate_slide(prs, 0)
        build_slide17(s17)

        print("Slide 18 (검색 예시) 구축...")
        s18 = duplicate_slide(prs, 0)
        build_slide18(s18)

        print("슬라이드 순서 재정렬 (기존 slide 2 제거)...")
        reorder_slides(prs, [0, 2, 3, 4, 5, 6, 7])

    prs.save(OUTPUT)
    print("\n저장 완료: " + OUTPUT)



if __name__ == "__main__":
    main()
