"""RAG 시스템 PPT 생성 스크립트
Usage: python scripts/create_rag_ppt.py
Output: RAG_시스템_소개.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── 색상 팔레트 ──────────────────────────────────────────────
CN    = RGBColor(0x1E, 0x3A, 0x5F)  # navy
CB    = RGBColor(0x2B, 0x7F, 0xD4)  # blue
CLB   = RGBColor(0xD4, 0xE9, 0xF7)  # light blue
CLLB  = RGBColor(0xEF, 0xF7, 0xFF)  # very light blue
COR   = RGBColor(0xE8, 0x6A, 0x2C)  # orange
CW    = RGBColor(0xFF, 0xFF, 0xFF)  # white
CBK   = RGBColor(0x1A, 0x1A, 0x2E)  # near black
CGR   = RGBColor(0xF2, 0xF2, 0xF2)  # gray
CDG   = RGBColor(0x88, 0x88, 0x99)  # dark gray
CTEAL = RGBColor(0x16, 0x8A, 0x7A)  # teal
CPP   = RGBColor(0x8E, 0x44, 0xAD)  # purple
CYEL  = RGBColor(0xFF, 0xF3, 0xCC)  # yellow bg
CYEL2 = RGBColor(0xFF, 0xCC, 0x00)  # yellow border
CYEL3 = RGBColor(0x7D, 0x5A, 0x00)  # yellow text
CDARK = RGBColor(0x2D, 0x2D, 0x2D)  # code bg
CCODE = RGBColor(0xF8, 0xF8, 0xF2)  # code text

W, H   = 13.33, 7.5
FONT   = "Malgun Gothic"
FMONO  = "Consolas"


# ─── 헬퍼 ────────────────────────────────────────────────────
def new_prs():
    p = Presentation()
    p.slide_width  = Inches(W)
    p.slide_height = Inches(H)
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg_color(slide, c=CW):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = c

def box(slide, x, y, w, h, fill=None, border=None, bw=0.75):
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(bw)
    else:
        shp.line.fill.background()
    return shp

def txt(slide, text, x, y, w, h, sz=12, bold=False, clr=CBK,
        align=PP_ALIGN.LEFT, fill=None, font=FONT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    if fill:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = clr
    r.font.name = font
    return tb

def multi_txt(slide, lines, x, y, w, h, sz=10, clr=CBK, font=FONT, sp=2):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if i > 0 and sp:
            p.space_before = Pt(sp)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(sz)
        r.font.color.rgb = clr
        r.font.name = font
    return tb

def title_bar(slide, title, sub=None):
    box(slide, 0, 0, W, 1.0, fill=CN)
    txt(slide, title, 0.35, 0.10, W-0.7, 0.62, sz=20, bold=True, clr=CW)
    if sub:
        txt(slide, sub, 0.35, 0.70, W-0.7, 0.27, sz=10,
            clr=RGBColor(0xB0, 0xCE, 0xEE))

def msg_box(slide, msg, y=1.06):
    box(slide, 0.35, y, W-0.7, 0.50, fill=CLB, border=CB, bw=0.5)
    txt(slide, f'"{msg}"', 0.50, y+0.04, W-1.0, 0.42,
        sz=11.5, clr=CN, italic=True)

def tbl(slide, headers, rows, x, y, w, h, col_ws=None,
        hfill=CN, sz=10, hsz=10):
    nc = len(headers)
    nr = len(rows) + 1
    t = slide.shapes.add_table(
        nr, nc, Inches(x), Inches(y), Inches(w), Inches(h)).table
    if col_ws:
        for i, cw in enumerate(col_ws):
            t.columns[i].width = Inches(cw)
    for j, h_text in enumerate(headers):
        c = t.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = hfill
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h_text
        r.font.bold = True; r.font.color.rgb = CW
        r.font.size = Pt(hsz); r.font.name = FONT
    for i, row in enumerate(rows):
        fc = CLLB if i % 2 == 0 else CW
        for j, val in enumerate(row):
            c = t.cell(i+1, j)
            c.fill.solid(); c.fill.fore_color.rgb = fc
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(sz); r.font.color.rgb = CBK; r.font.name = FONT
    return t

def code_box(slide, text, x, y, w, h, sz=9):
    box(slide, x, y, w, h, fill=CDARK)
    txt(slide, text, x+0.12, y+0.08, w-0.24, h-0.16,
        sz=sz, clr=CCODE, font=FMONO)


# ─── 슬라이드 생성 ────────────────────────────────────────────
def make_title_slide(prs):
    s = blank(prs)
    bg_color(s, CN)
    box(s, 0, 5.8, W, 0.09, fill=COR)
    box(s, 0, 5.92, W, 0.06, fill=CB)
    txt(s, "HAZOP 공정 안전", 1.0, 1.3, W-2.0, 1.1,
        sz=46, bold=True, clr=CW, align=PP_ALIGN.CENTER)
    txt(s, "RAG 시스템", 1.0, 2.35, W-2.0, 1.1,
        sz=46, bold=True, clr=CB, align=PP_ALIGN.CENTER)
    txt(s, "화공 도메인 특화  검색-증강 생성(RAG) 파이프라인",
        1.0, 3.7, W-2.0, 0.6,
        sz=16, clr=RGBColor(0xB0, 0xCE, 0xEE), align=PP_ALIGN.CENTER)
    db_labels = ["사고사례", "화학물질", "법령", "설계지침", "화공일반정보"]
    db_colors = [COR, CB, CTEAL, CPP, RGBColor(0x27, 0xAE, 0x60)]
    bw = 1.9; gap = 0.14
    sx = (W - (len(db_labels)*bw + (len(db_labels)-1)*gap)) / 2
    txt(s, "5종 Knowledge Base", sx, 4.4, len(db_labels)*(bw+gap), 0.3,
        sz=10, clr=CDG, align=PP_ALIGN.CENTER)
    for i, (label, color) in enumerate(zip(db_labels, db_colors)):
        bx = sx + i*(bw+gap)
        box(s, bx, 4.72, bw, 0.48, fill=color)
        txt(s, label, bx, 4.72, bw, 0.48,
            sz=13, bold=True, clr=CW, align=PP_ALIGN.CENTER)


def make_slide1(prs):
    s = blank(prs); bg_color(s)
    title_bar(s, "Slide 1  |  RAG 시스템 전체 구조")
    msg_box(s, "질문이 들어오면 → 의도를 파악하고 → 관련 문서를 찾고 → 그 내용을 근거로 LLM이 답변을 생성")

    steps = [
        ("① 질문 입력",  "질문\n+ P&ID JSON"),
        ("② 의도 분류",  "risk / law\ndesign / general"),
        ("③ 개체 추출",  "물질명\n설비명"),
        ("④ 1차 검색",   "Hybrid Search\nEmbedding+BM25"),
        ("⑤ 2차 검색",   "Re-ranking\nCross-Encoder"),
        ("⑥ Case 분류", "검색 결과\n패턴 판별"),
        ("⑦ 답변 생성",  "LLM\n(인용 포함)"),
    ]
    colors = [CN, CB, CB, RGBColor(0x16,0x6A,0xAA), RGBColor(0x0D,0x5A,0x9A), COR, CTEAL]
    bw = 1.55; gap = 0.12
    sx = (W - (len(steps)*bw + (len(steps)-1)*gap)) / 2
    by = 1.85; bh = 1.3

    for i, (label, sub) in enumerate(steps):
        bx = sx + i*(bw+gap)
        box(s, bx, by, bw, bh, fill=colors[i])
        txt(s, label, bx+0.05, by+0.08, bw-0.1, bh*0.45,
            sz=9.5, bold=True, clr=CW, align=PP_ALIGN.CENTER)
        txt(s, sub, bx+0.05, by+bh*0.48, bw-0.1, bh*0.5,
            sz=8, clr=RGBColor(0xCC,0xE4,0xFF), align=PP_ALIGN.CENTER)
        if i < len(steps)-1:
            txt(s, "▶", bx+bw, by+bh/2-0.2, gap+0.08, 0.38,
                sz=10, clr=CDG, align=PP_ALIGN.CENTER)

    box(s, 0.35, 3.35, W-0.7, 0.48, fill=CGR, border=CDG, bw=0.3)
    txt(s, "5종 KB:  사고사례(accidents)  |  화학물질(chemicals)  |  법령(laws)  |  설계지침(designs)  |  화공일반정보(basics)",
        0.5, 3.38, W-1.0, 0.42, sz=9.5, clr=CDG, align=PP_ALIGN.CENTER)

    details = [
        ("② 의도 분류기", "risk → 사고사례+화학물질 DB  /  law → 법령 DB  /  design → 설계지침 DB  /  general → 검색 없이 안내 반환"),
        ("③ 개체 추출기", "물질명·설비명 추출 → 사고사례 DB 필드 기준 정밀 검색에 활용"),
        ("⑥ Case 분류",  "물질·설비 일치 패턴으로 1-1~1-5, 2-1~2-4, 3-1~3-4 자동 결정 → 이후 프롬프트 전략 결정"),
    ]
    dy = 4.0
    for label, desc in details:
        box(s, 0.35, dy, 1.65, 0.4, fill=CLB)
        txt(s, label, 0.35, dy, 1.65, 0.4, sz=9.5, bold=True, clr=CN, align=PP_ALIGN.CENTER)
        txt(s, desc, 2.08, dy+0.03, W-2.45, 0.34, sz=9.5, clr=CBK)
        dy += 0.44


def make_slide2(prs):
    s = blank(prs); bg_color(s)
    title_bar(s, "Slide 2  |  지식 베이스 — Corpus 구성")
    msg_box(s, "5종의 전문 문서를 목적에 맞게 가공하여 지식 베이스를 구성")

    headers = ["구분", "accidents", "chemicals", "laws", "designs", "basics"]
    rows = [
        ["내용",     "국내외 화학\n사고 이력",   "화학물질 안전정보\n460종",     "산업안전보건법 등",            "KOSHA Guide 등",                    "화공 전공서\n일반지식"],
        ["원본",     "Excel 3종",               "PDF\n(물질별 페이지)",          "PDF",                         "PDF\n(카테고리 폴더)",              "PDF\n(논문·교재·보고서)"],
        ["청킹",     "행 1개\n= 청크 1개",       "페이지 단위",                  "조항(제N조) 단위",             "페이지 단위",                       "페이지 단위"],
        ["특이사항", "5개 필드\n구조화",          "물질명\n별도 인덱싱",          "전문 제거, 2000자\n서브청킹, 메타데이터",  "KOSHA/KGS 지침명\n정제, 섹션 메타데이터", "챕터\n메타데이터"],
    ]
    tbl(s, headers, rows, 0.35, 1.72, W-0.7, 2.75,
        col_ws=[1.25, 2.0, 2.0, 2.1, 2.25, 2.0], sz=9, hsz=10)

    txt(s, "사고사례 구조화 포맷", 0.35, 4.6, 5.0, 0.32, sz=11, bold=True, clr=CN)
    code_box(s,
        "[사고내용] 톨루엔 취급 중 반응기에서 폭발 발생\n"
        "[관련설비] 반응기(Reactor)\n"
        "[관련물질] 톨루엔(Toluene)\n"
        "[사고유형] 폭발\n"
        "[사고원인] 냉각수 공급 중단으로 인한 폭주 반응",
        0.35, 4.97, 5.3, 1.55, sz=9.5)
    txt(s, "→ 5개 필드 구조화로 물질·설비 기준 정밀 검색 가능",
        0.35, 6.57, 5.3, 0.35, sz=9.5, clr=CDG, italic=True)

    txt(s, "법령 청킹 상세", 5.9, 4.6, 7.1, 0.32, sz=11, bold=True, clr=CN)
    box(s, 5.9, 4.97, 7.1, 1.55, fill=CLLB, border=CB, bw=0.5)
    multi_txt(s, [
        "  • 다른 테이블: 페이지 단위 분할",
        "  • 법령만 조항(제N조) 단위로 분할 → 조항별 독립 검색",
        "  • 2000자 초과 조항: 조항 번호 반복 삽입하며 서브청킹",
        "    (임베딩 품질 보호)",
        "  • 각 청크에 법령명·조항 번호 메타데이터 저장",
        "    → 검색 결과에 '산업안전보건법 제38조' 형태로 출처 표시",
    ], 6.0, 5.03, 6.8, 1.44, sz=10, clr=CBK, sp=2)


def make_slide3(prs):
    s = blank(prs); bg_color(s)
    title_bar(s, "Slide 3  |  Vector DB 구조")
    msg_box(s, "문서를 숫자 벡터로 변환하여 저장하되, 테이블별로 목적에 맞는 추가 인덱스를 구성")

    infos = [
        ("DB",       "LanceDB\n로컬 파일 기반"),
        ("테이블",   "5개\naccidents / chemicals\nlaws / designs / basics"),
        ("임베딩 모델", "BAAI/bge-m3\n1024차원, 한/영 다국어"),
        ("유사도",   "코사인 유사도\n(벡터 정규화 후 내적)"),
    ]
    bw = 2.88; gap = 0.18
    sx = (W - (len(infos)*bw + (len(infos)-1)*gap)) / 2
    for i, (label, val) in enumerate(infos):
        bx = sx + i*(bw+gap)
        box(s, bx, 1.73, bw, 1.15, fill=CLB, border=CB, bw=0.5)
        txt(s, label, bx+0.1, 1.78, bw-0.2, 0.32, sz=10, bold=True, clr=CN, align=PP_ALIGN.CENTER)
        txt(s, val,   bx+0.1, 2.1,  bw-0.2, 0.74, sz=9.5, clr=CBK,  align=PP_ALIGN.CENTER)

    txt(s, "테이블별 벡터 컬럼", 0.35, 3.05, 5.0, 0.35, sz=13, bold=True, clr=CN)
    tbl(s,
        ["테이블", "text_vector", "추가 벡터 컬럼", "용도"],
        [
            ["accidents", "O", "material_vector, equipment_vector",   "물질·설비 필드 정밀 매칭"],
            ["chemicals", "O", "chemical_name_vector",                "물질명 전용 유사도 검색"],
            ["laws",      "O", "title_vector, article_vector",        "법령명 + 조항(제N조) 검색"],
            ["designs",   "O", "title_vector, section_vector",        "지침명(KOSHA/KGS) + 섹션 번호 검색"],
            ["basics",    "O", "title_vector, chapter_vector",        "문서 제목 + 챕터명 기준 검색"],
        ],
        0.35, 3.44, W-0.7, 3.6,
        col_ws=[1.8, 1.55, 3.3, 5.25], sz=10, hsz=10)


def make_slide4(prs):
    s = blank(prs); bg_color(s)
    title_bar(s, "Slide 4  |  Query 전처리 — 의도 분류기(Intent Router) & 개체 추출기(Entity Extractor)")
    msg_box(s, "검색 전에 LLM이 질문의 의도와 핵심어를 먼저 파악하여 검색 전략을 결정")

    # 예시 질문 박스
    box(s, 0.35, 1.7, W-0.7, 0.48,
        fill=RGBColor(0xFF, 0xF3, 0xCC), border=RGBColor(0xFF, 0xCC, 0x00), bw=0.5)
    txt(s, '예시:  "톨루엔을 취급하는 반응기의 공정위험성을 알려줘."',
        0.5, 1.73, W-1.0, 0.42, sz=12, bold=True,
        clr=RGBColor(0x7D, 0x5A, 0x00))

    # ① Intent Router
    txt(s, "① 의도 분류기 (Intent Router)", 0.35, 2.3, 6.0, 0.35, sz=12, bold=True, clr=CN)
    intents = [
        ("risk",    "공정위험성", CB,    "사고사례 DB + 화학물질 DB 동시 검색"),
        ("law",     "법규 준수",  CTEAL, "산업안전보건법 등 법령 DB 검색"),
        ("design",  "설계 오류",  CPP,   "KOSHA Guide, KGS 등 설계지침 DB 검색"),
        ("general", "범위 외",    CDG,   "검색 없이 안내 메시지 반환"),
    ]
    iy = 2.72
    for intent, meaning, color, desc in intents:
        box(s, 0.35, iy, 1.05, 0.38, fill=color)
        txt(s, intent, 0.35, iy, 1.05, 0.38, sz=11, bold=True, clr=CW, align=PP_ALIGN.CENTER)
        txt(s, meaning, 1.47, iy+0.03, 1.25, 0.32, sz=10, bold=True, clr=CBK)
        txt(s, desc,    2.77, iy+0.03, 4.0,  0.32, sz=10, clr=CBK)
        iy += 0.43

    box(s, 0.35, 4.52, 6.5, 0.38, fill=CGR, border=CDG, bw=0.3)
    txt(s, '출력:  {"intents": ["risk"]}   •   복수 의도 감지 → 재질문 안내 반환',
        0.5, 4.55, 6.2, 0.32, sz=10, clr=CDG)

    # ② Entity Extractor
    txt(s, "② 개체 추출기 (Entity Extractor)", 7.1, 2.3, 5.9, 0.35, sz=12, bold=True, clr=CN)
    box(s, 7.1, 2.72, 5.9, 1.5, fill=CLLB, border=CB, bw=0.5)
    multi_txt(s, [
        "  • 질문에서 물질명·설비명 추출",
        '  • 출력: {"target_material": "톨루엔",',
        '           "target_equipment": "반응기"}',
        "  • 추출 실패(None) → 전체 텍스트 검색으로 대체",
        "  • 추출 값은 사고사례 DB의 material / equipment",
        "    필드 검색에 직접 사용",
    ], 7.2, 2.78, 5.7, 1.4, sz=10, clr=CBK, sp=2)

    txt(s, "추출값 활용처", 7.1, 4.32, 5.9, 0.32, sz=11, bold=True, clr=CN)
    uses = [("사고사례 물질 검색", CB), ("사고사례 설비 검색", CB), ("P&ID 프롬프트 삽입", COR)]
    ux = 7.1
    for label, color in uses:
        box(s, ux, 4.68, 1.88, 0.46, fill=color)
        txt(s, label, ux, 4.68, 1.88, 0.46, sz=9.5, bold=True, clr=CW, align=PP_ALIGN.CENTER)
        ux += 1.98

    box(s, 0.35, 5.1, W-0.7, 0.4, fill=CLB, border=CB, bw=0.3)
    txt(s, "※ 분류 결과에 따라 어떤 DB를 검색할지, 몇 건을 가져올지, 어떤 프롬프트를 쓸지가 달라진다",
        0.5, 5.13, W-1.0, 0.34, sz=10, bold=True, clr=CN)


def make_slide5(prs):
    s = blank(prs); bg_color(s)
    title_bar(s, "Slide 5  |  1차 검색 — Hybrid Search")
    msg_box(s, "의미 기반 검색과 키워드 기반 검색을 결합하여 두 방식의 장점을 동시에 활용")

    # 좌: 비교 표
    txt(s, "두 가지 검색 방식", 0.35, 1.72, 6.5, 0.32, sz=12, bold=True, clr=CN)
    tbl(s,
        ["", "Semantic Search\n(Embedding)", "Lexical Search\n(BM25)"],
        [
            ["방식",   "문장 전체 의미를 벡터로 표현\n유사 벡터 검색",  "단어 빈도 기반 관련도 계산"],
            ["강점",   "표현이 달라도\n의미 유사한 문서 탐색",          "전문 용어·법조항 번호 정확 매칭"],
            ["모델",   "BAAI/bge-m3 (1024차원)",                       "BM25 (LanceDB FTS)"],
        ],
        0.35, 2.07, 6.5, 2.45, col_ws=[1.1, 2.7, 2.7], sz=10, hsz=10)

    txt(s, "Hybrid Score", 0.35, 4.65, 6.5, 0.32, sz=12, bold=True, clr=CN)
    box(s, 0.35, 5.0, 6.5, 0.62, fill=CLB, border=CB, bw=0.5)
    txt(s, "score = 0.5 × semantic_score + 0.5 × BM25_score",
        0.5, 5.05, 6.2, 0.35, sz=13, bold=True, clr=CN,
        align=PP_ALIGN.CENTER, font=FMONO)
    txt(s, "BM25 점수는 최대값 기준 0~1 정규화 후 합산",
        0.5, 5.4, 6.2, 0.2, sz=9, clr=CDG, align=PP_ALIGN.CENTER)

    txt(s, "BM25 직관적 이해 (비전문가용)", 0.35, 5.72, 6.5, 0.28, sz=10, bold=True, clr=CBK)
    box(s, 0.35, 6.04, 6.5, 0.95, fill=CGR, border=CDG, bw=0.3)
    txt(s,
        '"이 단어가 이 문서에 자주 나오면서, 다른 문서들엔 별로 안 나올수록\n'
        '→ 핵심 키워드일 가능성이 높다."\n'
        '단, 문서가 길수록 자연히 더 많이 등장하므로 문서 길이로 나눠 공평하게 비교.',
        0.5, 6.08, 6.1, 0.87, sz=9.5, clr=CBK, italic=True)

    # 우: 한도 표
    txt(s, "테이블별 1차 검색 한도", 7.1, 1.72, 6.1, 0.32, sz=12, bold=True, clr=CN)
    tbl(s,
        ["테이블", "한도", "이유"],
        [
            ["accidents", "300건", "이후 필드 매칭으로 필터링 → 대량 조회 필요"],
            ["chemicals", "3건",   "물질명 직접 매칭, 소량으로 충분"],
            ["laws",      "10건",  "조항 단위 청크, 소수로 정확도 확보"],
            ["designs",   "10건",  "동일"],
            ["basics",    "10건",  "designs와 동일 (미구현, 예정)"],
        ],
        7.1, 2.07, 6.1, 2.45, col_ws=[1.5, 0.9, 3.7], sz=10, hsz=10)

    txt(s, "BM25 수식", 7.1, 4.65, 6.1, 0.32, sz=12, bold=True, clr=CN)
    code_box(s,
        "BM25(q,d) = Σ IDF(t) ·  TF(t,d)·(k₁+1)\n"
        "            ─────────────────────────────────────\n"
        "            TF(t,d) + k₁·(1 − b + b·|d|/avgdl)\n\n"
        "IDF: 희귀 단어에 높은 가중치  |  b: 문서 길이 정규화",
        7.1, 5.0, 6.1, 2.0, sz=9)


def make_slide6(prs):
    s = blank(prs); bg_color(s)
    title_bar(s, "Slide 6  |  사고사례 전용 검색 파이프라인")
    msg_box(s, "사고사례는 물질명·설비명 필드를 단계적으로 검색하여 일치 패턴을 자동 분류")

    txt(s, "4단계 검색 흐름", 0.35, 1.72, 6.4, 0.32, sz=12, bold=True, clr=CN)
    steps6 = [
        ("STEP 1", "필드별 Hybrid 검색",
         "material 필드 ← target_material → _material_relevance_score\n"
         "equipment 필드 ← target_equipment → _equipment_relevance_score\n"
         "두 결과를 ID 기준 병합 (score: max 취합)"),
        ("STEP 2", "임계치 필터링 (threshold = 0.7)",
         "both_res:     material ≥ 0.7  AND  equipment ≥ 0.7\n"
         "material_res: material ≥ 0.7\n"
         "equipment_res: equipment ≥ 0.7"),
        ("STEP 3", "Case별 Re-ranking",
         "1-1: material + equipment rerank 점수 평균\n"
         "1-2: material rerank  /  1-3: equipment rerank\n"
         "rerank score ≥ 0.7 통과, 최대 6건 반환"),
        ("STEP 4", "Case Code 자동 결정",
         "우선순위: 1-1 > 1-2 > 1-3 > 1-4 > 1-5"),
    ]
    scolors = [CB, RGBColor(0x16,0x6A,0xAA), RGBColor(0x0D,0x5A,0x9A), CN]
    sy = 2.1
    for i, (step, title6, desc) in enumerate(steps6):
        box(s, 0.35, sy, 1.0, 0.88, fill=scolors[i])
        txt(s, step, 0.35, sy+0.1, 1.0, 0.65, sz=9, bold=True, clr=CW, align=PP_ALIGN.CENTER)
        box(s, 1.42, sy, 5.25, 0.88, fill=CLLB, border=CLB, bw=0.3)
        txt(s, title6, 1.52, sy+0.03, 5.05, 0.28, sz=10, bold=True, clr=CN)
        txt(s, desc,   1.52, sy+0.3,  5.05, 0.55, sz=8.5, clr=CBK)
        if i < len(steps6)-1:
            txt(s, "▼", 0.6, sy+0.88, 0.7, 0.2, sz=9, clr=CDG, align=PP_ALIGN.CENTER)
        sy += 1.1

    txt(s, "Case Code", 7.0, 1.72, 6.2, 0.32, sz=12, bold=True, clr=CN)
    tbl(s,
        ["Code", "조건", "의미"],
        [
            ["1-1", "물질 ✓ + 설비 ✓",         "동일 물질·유사 설비 사고 이력 — 위험성 가장 높음"],
            ["1-2", "물질 ✓ + 설비 ✗",         "동일 물질 사고만 존재"],
            ["1-3", "물질 ✗ + 설비 ✓",         "유사 설비 사고만 존재"],
            ["1-4", "둘 다 ✗ + 화학물질 정보 ✓", "사고 이력 없음, 물질 안전정보로 대응"],
            ["1-5", "근거 없음",                 "답변 불가"],
        ],
        7.0, 2.1, 6.2, 2.85, col_ws=[0.7, 2.4, 3.1], sz=10, hsz=10)

    box(s, 7.0, 5.08, 6.2, 0.55, fill=CLB, border=CB, bw=0.5)
    txt(s, "이 Case Code가 이후 프롬프트 전략과 고정 문구(Intro/Outro)를 결정",
        7.15, 5.12, 5.9, 0.46, sz=10, bold=True, clr=CN)
    txt(s, "Law / Design도 동일 방식으로 분류 (2-1~2-4, 3-1~3-4)",
        7.0, 5.73, 6.2, 0.32, sz=10, clr=CBK, italic=True)


def make_slide7(prs):
    s = blank(prs); bg_color(s)
    title_bar(s, "Slide 7  |  2차 검색 — Re-ranking")
    msg_box(s, "1차 검색으로 확보한 후보를 Cross-Encoder 모델이 질문과 쌍으로 비교하여 관련도를 재계산")

    txt(s, "Bi-Encoder vs Cross-Encoder", 0.35, 1.72, 9.1, 0.32, sz=12, bold=True, clr=CN)
    tbl(s,
        ["", "Bi-Encoder (1차 임베딩)", "Cross-Encoder (2차 Re-rank)"],
        [
            ["방식",   "질문/문서 각각 인코딩 후 벡터 유사도 비교",   "질문+문서를 한 쌍으로 함께 입력"],
            ["속도",   "빠름 (벡터 캐시 가능)",                       "느림 (쌍마다 추론 필요)"],
            ["정확도", "상대적으로 낮음",                              "높음"],
            ["역할",   "전체 DB에서 후보 대량 추출",                   "소수 후보를 정밀 재정렬"],
        ],
        0.35, 2.07, 9.1, 2.95, col_ws=[1.4, 3.85, 3.85], sz=10, hsz=10)

    box(s, 0.35, 5.12, 9.1, 0.42, fill=CGR, border=CDG, bw=0.3)
    txt(s, "모델:  BAAI/bge-reranker-v2-m3   •   속도가 느리므로 1차 검색으로 후보를 줄인 뒤 적용 (2-stage retrieval)",
        0.5, 5.15, 8.8, 0.36, sz=10, clr=CBK)

    txt(s, "테이블별 Re-rank 설정", 9.65, 1.72, 3.5, 0.32, sz=12, bold=True, clr=CN)
    tbl(s,
        ["테이블", "최종 반환", "임계치", "비고"],
        [
            ["accidents", "6건", "0.7", "물질/설비 필드 기준"],
            ["chemicals", "1건", "0.7", "물질명 필드 기준"],
            ["laws",      "3건", "0.5", "법조문 특성상 완화"],
            ["designs",   "3건", "0.5", "동일"],
        ],
        9.65, 2.07, 3.5, 2.3, col_ws=[1.35, 0.8, 0.8, 0.55], sz=9.5, hsz=10)

    box(s, 9.65, 4.48, 3.5, 0.65, fill=CLLB, border=CB, bw=0.3)
    txt(s, "laws·designs 임계치 0.5로 완화:\n법조문은 질문 표현과 달라도 내용적으로\n관련 있는 경우가 많음",
        9.75, 4.52, 3.3, 0.58, sz=9, clr=CBK)

    txt(s, "2-stage Retrieval 흐름", 0.35, 5.68, 9.1, 0.3, sz=11, bold=True, clr=CN)
    stages = [
        ("전체 DB\n(수천~수만 청크)", CGR, CBK),
        ("1차 검색\n(Hybrid Search)\n→ 후보 추출", CB, CW),
        ("2차 검색\n(Cross-Encoder\nRe-ranking)", CN, CW),
        ("최종 결과\n(3~6건)", COR, CW),
    ]
    sw = 2.1; sx = 0.35; sh = 0.97; sy7 = 6.05
    for i, (label, fc, tc) in enumerate(stages):
        bx = sx + i*(sw+0.32)
        box(s, bx, sy7, sw, sh, fill=fc, border=CDG if fc == CGR else None)
        txt(s, label, bx, sy7, sw, sh, sz=9, bold=True, clr=tc, align=PP_ALIGN.CENTER)
        if i < len(stages)-1:
            txt(s, "→", bx+sw, sy7+sh/2-0.2, 0.35, 0.38, sz=14, clr=CDG, align=PP_ALIGN.CENTER)


def make_slide8(prs):
    s = blank(prs); bg_color(s)
    title_bar(s, "Slide 8  |  프롬프트 구성 & 답변 생성")
    msg_box(s, "검색된 문서에 인용 번호를 붙이고, P&ID 정보와 함께 조합하여 LLM이 근거 있는 답변을 생성")

    txt(s, "문서 → Context 포맷 ([CITE_N])", 0.35, 1.72, 5.9, 0.32, sz=12, bold=True, clr=CN)
    code_box(s,
        "[CITE_1]\n"
        "문서 제목: domestic_accident.xlsx\n"
        "페이지: -  |  유사도: 0.87\n"
        "본문:\n"
        "[사고내용] 톨루엔 취급 중 반응기에서 폭발...\n"
        "[관련물질] 톨루엔  [관련설비] 반응기",
        0.35, 2.07, 5.9, 1.52, sz=9.5)
    txt(s, "→ 각 문서에 [CITE_N] 번호 부여, LLM이 답변 내 인용 표시\n→ 최종 API 응답에 출처 목록(문서명·PDF 링크·페이지) 포함",
        0.35, 3.65, 5.9, 0.55, sz=9.5, clr=CBK)

    txt(s, "P&ID (facility_info) 활용", 0.35, 4.3, 5.9, 0.32, sz=12, bold=True, clr=CN)
    box(s, 0.35, 4.66, 5.9, 1.75, fill=CLLB, border=CB, bw=0.4)
    multi_txt(s, [
        "  • 설비 도면 JSON을 함께 전달 (선택)",
        "  • JSON을 그대로 넣지 않고, 한국어 텍스트로 변환 후 프롬프트에 삽입",
        "  • 포함 정보: 설비 ID·명칭·설계 사양·취급 물질",
        "               + 연결 배관 흐름 + 설치 구성품(밸브·스트레이너 등)",
        "  • 입력이 있을수록 법규·설계 검토 답변 품질 향상",
    ], 0.45, 4.72, 5.7, 1.65, sz=10, clr=CBK, sp=2)

    txt(s, "Intent별 LLM 호출 구조", 6.45, 1.72, 6.7, 0.32, sz=12, bold=True, clr=CN)
    tbl(s,
        ["Intent", "호출 방식", "내용"],
        [
            ["risk (Case 1-1~1-3)", "병렬 2회",  "호출A: 위험 특성+대응 (2문장)\n호출B: 사고 사례 목록"],
            ["risk (Case 1-4)",     "1회",        "화학물질 안전정보 기반 설명"],
            ["law / design",        "순차 2회",   "호출1: Case Selector (위반/준수/불명확)\n→ 호출2: 본문 생성"],
        ],
        6.45, 2.07, 6.7, 2.35, col_ws=[2.0, 1.3, 3.4], sz=10, hsz=10)

    txt(s, "고정 Intro / Outro", 6.45, 4.52, 6.7, 0.32, sz=12, bold=True, clr=CN)
    box(s, 6.45, 4.88, 6.7, 1.2, fill=CLLB, border=CB, bw=0.4)
    multi_txt(s, [
        "  • law / design은 Case 코드에 따라 고정 도입/마무리 문구 삽입",
        "  • LLM 본문을 고정 문구로 감싸는 구조 → 일관된 어조 유지",
        '  • 예) Law 위반 케이스:',
        '    "검색된 문서를 기준으로 검토한 결과,',
        '     법적 요구사항을 충족하지 않는 것으로 확인되었습니다."',
    ], 6.55, 4.93, 6.5, 1.12, sz=9.5, clr=CBK, sp=1)

    txt(s, "프롬프트 규칙 예시 (risk 1-1)", 6.45, 6.2, 6.7, 0.3, sz=10, bold=True, clr=CN)
    code_box(s,
        "[규칙]\n"
        "• [근거 문서] 밖의 사실은 쓰지 마십시오.\n"
        "• 첫 문장: 위험적 특성 50자 내\n"
        "• 둘째 문장: 대응/제안 50자 내",
        6.45, 6.55, 6.7, 0.9, sz=9)


# ─── 실행 ─────────────────────────────────────────────────────
prs = new_prs()
make_title_slide(prs)
make_slide1(prs)
make_slide2(prs)
make_slide3(prs)
make_slide4(prs)
make_slide5(prs)
make_slide6(prs)
make_slide7(prs)
make_slide8(prs)

out = "RAG_시스템_소개.pptx"
prs.save(out)
print(f"저장 완료: {out}")
